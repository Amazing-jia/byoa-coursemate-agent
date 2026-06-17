import sys
from pathlib import Path

VENDOR_DIR = Path(__file__).resolve().parent / ".vendor"
if VENDOR_DIR.exists():
    sys.path.insert(0, str(VENDOR_DIR))

from pptx import Presentation
from PyPDF2 import PdfReader


def read_pptx(path):
    prs = Presentation(path)
    texts = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

        if slide_texts:
            texts.append(f"[Slide {slide_index}]\n" + "\n".join(slide_texts))

    return "\n\n".join(texts)


def read_pdf(path):
    reader = PdfReader(path)
    texts = []

    for page_index, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text and text.strip():
            texts.append(f"[Page {page_index}]\n{text.strip()}")

    return "\n\n".join(texts)


def read_txt(path):
    with open(path, "r", encoding="utf-8") as file:
        return file.read()


def read_local_file(path):
    file_path = Path(path)

    if not file_path.exists():
        return f"File not found: {file_path}"

    suffix = file_path.suffix.lower()
    if suffix == ".pptx":
        return read_pptx(file_path)
    if suffix == ".pdf":
        return read_pdf(file_path)
    if suffix == ".txt":
        return read_txt(file_path)

    return f"Unsupported file format: {suffix}"


def search_in_text(text, keyword, max_results=20):
    if not keyword.strip():
        return "Please provide a keyword to search."

    blocks = [block.strip() for block in text.split("\n\n") if block.strip()]
    keyword_lower = keyword.lower()
    results = []

    for block in blocks:
        if keyword_lower in block.lower():
            results.append(block)

    if not results:
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        results = [line for line in lines if keyword_lower in line.lower()]

    return "\n\n---\n\n".join(results[:max_results]) if results else "No relevant content found."


def summarize_text(text, max_chars=1500):
    clean_lines = [line.strip() for line in text.splitlines() if line.strip()]
    summary_source = " ".join(clean_lines)
    if len(summary_source) <= max_chars:
        return summary_source
    return summary_source[:max_chars].rstrip() + "..."
